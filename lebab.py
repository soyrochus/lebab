import argparse
import asyncio
import os
import shutil
import tempfile

from dotenv import load_dotenv
from langchain_openai import ChatOpenAI# Requires langchain>=0.3
import docx  # python-docx for Word files
from pptx import Presentation  # python-pptx for PowerPoint files

# Estimate a maximum character count per translation request.
# This is a high-level approximation based on ~10K tokens.
MAX_CHUNK_SIZE = 10000

def init_llm():
    """
    Initialize the LLM using ChatOpenAI.
    
    Design rationale:
    - Encapsulates all OpenAI-specific initialization.
    - Uses python-dotenv to load the API key.
    - Allows for later substitution with a different initializer.
    """
    load_dotenv()  # load environment variables from .env file
    openai_api_key = os.getenv("OPENAI_API_KEY")
    if not openai_api_key:
        raise ValueError("OPENAI_API_KEY not found in environment variables")
    
    # Create the ChatOpenAI instance.
    # Temperature is set to 0 for deterministic translation.
    llm = ChatOpenAI(openai_api_key=openai_api_key, temperature=0, model_name="gpt-3.5-turbo")
    return llm

async def translate_text_async(text, input_lang, output_lang, llm):
    """
    Translate the given text asynchronously using the LLM.
    
    Design rationale:
    - Constructs a prompt instructing the LLM to translate while preserving formatting.
    - Uses the async ainvoke method so that the translation call does not block.
    """
    prompt = (
        f"Translate the following text from {input_lang} to {output_lang} "
        "while preserving formatting and document structure:\n\n" + text
    )
    try:
        # Using ainvoke (async call) with a simple message structure.
        response = await llm.ainvoke([{"role": "user", "content": prompt}])
        # We assume the response object has a 'content' attribute.
        translated_text = response.content if hasattr(response, "content") else response
        return translated_text
    except Exception as e:
        print(f"Error during translation: {e}")
        # On error, return the original text as a fallback.
        return text

class DocumentTranslator:
    """
    Base class for document translators.
    
    Design rationale:
    - Provides a common interface for reading, updating, and writing documents.
    - Makes it easy to extend support for new formats (e.g., Excel).
    """
    def __init__(self, file_path):
        self.file_path = file_path
        self.blocks = []  # Each block represents a text unit from the document

    def read_document(self):
        raise NotImplementedError

    def write_document(self, target_path):
        raise NotImplementedError

    def update_blocks(self, translated_blocks):
        """
        Update the internal document model with the translated text.
        Must be implemented by subclasses.
        """
        raise NotImplementedError

class DocxTranslator(DocumentTranslator):
    """
    Translator for Word (.docx) documents using python-docx.
    
    Design rationale:
    - Iterates over paragraphs as blocks.
    - Keeps track of block indices so that after translation,
      we can replace the original text accurately.
    """
    def __init__(self, file_path):
        super().__init__(file_path)
        self.doc = None

    def read_document(self):
        try:
            self.doc = docx.Document(self.file_path)
            self.blocks = []
            # Extract each paragraph as a separate block.
            for i, para in enumerate(self.doc.paragraphs):
                self.blocks.append({"type": "paragraph", "index": i, "text": para.text})
        except Exception as e:
            print(f"Error reading DOCX file: {e}")

    def update_blocks(self, translated_blocks):
        # Replace paragraph texts with their translated versions.
        for block in translated_blocks:
            if block["type"] == "paragraph":
                try:
                    self.doc.paragraphs[block["index"]].text = block["translated_text"]
                except Exception as e:
                    print(f"Error updating paragraph {block['index']}: {e}")

    def write_document(self, target_path):
        try:
            self.doc.save(target_path)
        except Exception as e:
            print(f"Error writing DOCX file: {e}")

class PptxTranslator(DocumentTranslator):
    """
    Translator for PowerPoint (.pptx) presentations using python-pptx.
    
    Design rationale:
    - Iterates over slides and text-containing shapes.
    - Stores slide and shape indices to accurately replace text after translation.
    """
    def __init__(self, file_path):
        super().__init__(file_path)
        self.prs = None

    def read_document(self):
        try:
            self.prs = Presentation(self.file_path)
            self.blocks = []
            # Iterate through slides and shapes that contain text.
            for slide_index, slide in enumerate(self.prs.slides):
                for shape_index, shape in enumerate(slide.shapes):
                    if hasattr(shape, "text") and shape.text:
                        self.blocks.append({
                            "type": "shape",
                            "slide_index": slide_index,
                            "shape_index": shape_index,
                            "text": shape.text
                        })
        except Exception as e:
            print(f"Error reading PPTX file: {e}")

    def update_blocks(self, translated_blocks):
        # Update each shape with the translated text.
        for block in translated_blocks:
            if block["type"] == "shape":
                try:
                    slide = self.prs.slides[block["slide_index"]]
                    shape = slide.shapes[block["shape_index"]]
                    shape.text = block["translated_text"]
                except Exception as e:
                    print(f"Error updating slide {block['slide_index']} shape {block['shape_index']}: {e}")

    def write_document(self, target_path):
        try:
            self.prs.save(target_path)
        except Exception as e:
            print(f"Error writing PPTX file: {e}")

async def process_translation(translator, input_lang, output_lang, llm):
    """
    Process the translation of document blocks in manageable chunks.
    
    Design rationale:
    - Groups text blocks until a maximum estimated character size is reached.
    - Uses a delimiter ("---") to join blocks so that after translation,
      we can split the translated text back into individual blocks.
    - Ensures that blocks are never split, preserving document structure.
    - Catches errors during translation and proceeds with the rest.
    """
    translator.read_document()
    if not translator.blocks:
        print("No text blocks found for translation.")
        return

    translated_blocks = []
    current_chunk = []
    current_chunk_size = 0

    for block in translator.blocks:
        block_text = block["text"]
        block_size = len(block_text)
        # If adding this block would exceed the max size, translate the current chunk.
        if current_chunk and (current_chunk_size + block_size > MAX_CHUNK_SIZE):
            chunk_text = "\n---\n".join(b["text"] for b in current_chunk)
            translated_chunk = await translate_text_async(chunk_text, input_lang, output_lang, llm)
            # Split by the delimiter to map back to each block.
            translated_parts = translated_chunk.split("\n---\n")
            if len(translated_parts) != len(current_chunk):
                print("Warning: Mismatch in block count after translation. Some blocks may be left unmodified.")
                for i, b in enumerate(current_chunk):
                    b["translated_text"] = translated_parts[i] if i < len(translated_parts) else b["text"]
            else:
                for b, t in zip(current_chunk, translated_parts):
                    b["translated_text"] = t
            translated_blocks.extend(current_chunk)
            current_chunk = []
            current_chunk_size = 0

        current_chunk.append(block)
        current_chunk_size += block_size

    # Translate any remaining blocks.
    if current_chunk:
        chunk_text = "\n---\n".join(b["text"] for b in current_chunk)
        translated_chunk = await translate_text_async(chunk_text, input_lang, output_lang, llm)
        translated_parts = translated_chunk.split("\n---\n")
        if len(translated_parts) != len(current_chunk):
            print("Warning: Mismatch in block count after translation. Some blocks may be left unmodified.")
            for i, b in enumerate(current_chunk):
                b["translated_text"] = translated_parts[i] if i < len(translated_parts) else b["text"]
        else:
            for b, t in zip(current_chunk, translated_parts):
                b["translated_text"] = t
        translated_blocks.extend(current_chunk)

    # Update the document model with the translated text.
    translator.update_blocks(translated_blocks)

def construct_target_filename(input_file, output_lang):
    """
    Construct a target filename by appending the output language code.
    
    Example: 'MyText.docx' becomes 'MyText_EN.docx' for output_lang='EN'.
    """
    base, ext = os.path.splitext(input_file)
    return f"{base}_{output_lang}{ext}"

async def main():
    parser = argparse.ArgumentParser(description="Lebab: Document Translator Preserving Structure")
    parser.add_argument("input_file", help="Path to the input document")
    parser.add_argument("input_lang", help="Input language code (e.g., ES)")
    parser.add_argument("output_lang", help="Output language code (e.g., EN)")
    parser.add_argument("-t", "--to-file", help="Target file path")
    args = parser.parse_args()

    input_file = args.input_file
    target_file = args.to_file if args.to_file else construct_target_filename(input_file, args.output_lang)

    # Step 1: Create a temporary copy of the file.
    try:
        temp_dir = tempfile.gettempdir()
        temp_file = os.path.join(temp_dir, os.path.basename(input_file))
        shutil.copy(input_file, temp_file)
        print(f"Temporary file created at {temp_file}")
    except Exception as e:
        print(f"Error creating temporary file: {e}")
        return

    # Select the appropriate translator based on file extension.
    ext = os.path.splitext(input_file)[1].lower()
    if ext == ".docx":
        translator = DocxTranslator(temp_file)
    elif ext == ".pptx":
        translator = PptxTranslator(temp_file)
    else:
        print("Unsupported file format. Only .docx and .pptx are supported.")
        return

    # Step 2: Initialize the LLM.
    try:
        llm = init_llm()
    except Exception as e:
        print(f"Error initializing LLM: {e}")
        return

    # Step 3: Process translation (chunked to avoid excessive text per request).
    await process_translation(translator, args.input_lang, args.output_lang, llm)

    # Step 4: Write out the translated document.
    translator.write_document(target_file)
    print(f"Translated document saved to {target_file}")

if __name__ == "__main__":
    asyncio.run(main())
